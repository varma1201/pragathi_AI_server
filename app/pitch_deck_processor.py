"""
Pitch Deck Processor Module
Extracts idea information from PDF and PowerPoint pitch decks
"""

import os
import re
from typing import Dict, Optional, Tuple
from pathlib import Path

# PDF processing
try:
    from pypdf import PdfReader
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("Warning: PDF processing not available. Install pypdf and pdfplumber.")

# PowerPoint processing
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: PowerPoint processing not available. Install python-pptx.")

# AI-based text extraction and summarization
from langchain_openai import ChatOpenAI
from dotenv import load_dotenv

load_dotenv()


class PitchDeckProcessor:
    """Process PDF and PowerPoint pitch decks to extract startup idea information"""
    
    def __init__(self):
        """Initialize the pitch deck processor"""
        self.openai_api_key = os.getenv("OPENAI_API_KEY")
        if self.openai_api_key:
            self.llm = ChatOpenAI(
                openai_api_key=self.openai_api_key,
                temperature=0.2,
                model="gpt-4o-mini",  # Using mini for cost efficiency
                max_tokens=2000
            )
        else:
            self.llm = None
            print("Warning: OpenAI API key not found. AI extraction will not be available.")
    
    def process_pitch_deck(self, file_path: str) -> Dict[str, str]:
        """
        Process a pitch deck file (PDF or PPT) and extract idea information
        
        Args:
            file_path: Path to the pitch deck file
            
        Returns:
            Dictionary with 'idea_name' and 'idea_concept'
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"Pitch deck file not found: {file_path}")
        
        file_extension = file_path.suffix.lower()
        
        # Extract text based on file type
        if file_extension == '.pdf':
            raw_text = self._extract_from_pdf(str(file_path))
        elif file_extension in ['.ppt', '.pptx']:
            raw_text = self._extract_from_pptx(str(file_path))
        else:
            raise ValueError(f"Unsupported file format: {file_extension}. Supported: .pdf, .ppt, .pptx")
        
        # Extract structured information using AI
        if self.llm and raw_text:
            return self._extract_idea_with_ai(raw_text)
        else:
            # Fallback to basic extraction
            return self._extract_idea_basic(raw_text)
    
    def _extract_from_pdf(self, pdf_path: str) -> str:
        """Extract text from PDF file"""
        if not PDF_AVAILABLE:
            raise ImportError("PDF processing libraries not available. Install pypdf and pdfplumber.")
        
        text_content = []
        
        try:
            # Try with pdfplumber first (better for formatted PDFs)
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(f"\n--- Page {page_num} ---\n{page_text}")
        except Exception as e:
            print(f"pdfplumber failed: {e}. Trying pypdf...")
            
            # Fallback to pypdf
            try:
                reader = PdfReader(pdf_path)
                for page_num, page in enumerate(reader.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(f"\n--- Page {page_num} ---\n{page_text}")
            except Exception as e2:
                raise Exception(f"Failed to extract PDF text: {e2}")
        
        return "\n".join(text_content)
    
    def _extract_from_pptx(self, pptx_path: str) -> str:
        """Extract text from PowerPoint file"""
        if not PPTX_AVAILABLE:
            raise ImportError("PowerPoint processing library not available. Install python-pptx.")
        
        text_content = []
        
        try:
            prs = Presentation(pptx_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"\n--- Slide {slide_num} ---"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells if cell.text])
                            if row_text:
                                slide_text.append(row_text)
                
                text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
            
        except Exception as e:
            raise Exception(f"Failed to extract PowerPoint text: {e}")
    
    def _extract_idea_with_ai(self, raw_text: str) -> Dict[str, str]:
        """
        Use AI to extract structured idea information from raw text
        
        Args:
            raw_text: Raw text extracted from pitch deck
            
        Returns:
            Dictionary with 'idea_name' and 'idea_concept'
        """
        prompt = f"""
You are analyzing a startup pitch deck. Extract the following information from the text:

1. **Idea Name**: The name/title of the startup or product (be concise, 1-10 words)
2. **Idea Concept**: A comprehensive description of what the startup does, including:
   - Problem being solved
   - Solution/product description
   - Target market
   - Key features or unique value proposition
   - Business model (if mentioned)
   - Any specific focus on geography (especially India) or industry

PITCH DECK TEXT:
{raw_text[:4000]}  # Limit to first 4000 chars to stay within token limits

Provide your response in this exact JSON format:
{{
    "idea_name": "The startup/product name",
    "idea_concept": "A detailed 2-4 paragraph description covering problem, solution, target market, key features, and business model"
}}

Be thorough in the idea_concept - include all relevant details that would help evaluate the startup.
"""
        
        try:
            response = self.llm.invoke(prompt)
            response_text = response.content
            
            # Extract JSON from response
            import json
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            if json_match:
                extracted_data = json.loads(json_match.group())
                
                return {
                    "idea_name": extracted_data.get("idea_name", "Unnamed Startup"),
                    "idea_concept": extracted_data.get("idea_concept", raw_text[:1000])
                }
            else:
                # Fallback if JSON parsing fails
                return self._extract_idea_basic(raw_text)
                
        except Exception as e:
            print(f"AI extraction failed: {e}. Using basic extraction.")
            return self._extract_idea_basic(raw_text)
    
    def _extract_idea_basic(self, raw_text: str) -> Dict[str, str]:
        """
        Basic extraction without AI (fallback method)
        
        Args:
            raw_text: Raw text extracted from pitch deck
            
        Returns:
            Dictionary with 'idea_name' and 'idea_concept'
        """
        lines = raw_text.split('\n')
        lines = [line.strip() for line in lines if line.strip()]
        
        # Try to find the idea name (usually in first few slides/pages)
        idea_name = "Startup Idea"
        for line in lines[:20]:  # Check first 20 lines
            # Look for title-like patterns
            if len(line) < 100 and len(line.split()) <= 10:
                # Avoid common headers
                if not any(x in line.lower() for x in ['slide', 'page', 'confidential', 'proprietary']):
                    idea_name = line
                    break
        
        # Use first 1500 characters as concept
        idea_concept = " ".join(lines[:50])[:1500]
        
        return {
            "idea_name": idea_name,
            "idea_concept": idea_concept
        }
    
    def process_and_validate(self, file_path: str, custom_weights: Optional[Dict[str, float]] = None) -> Tuple[Dict[str, str], Dict]:
        """
        Process pitch deck and prepare for validation
        
        Args:
            file_path: Path to pitch deck file
            custom_weights: Optional custom cluster weights
            
        Returns:
            Tuple of (extracted_info, validation_payload)
        """
        # Extract idea information
        extracted_info = self.process_pitch_deck(file_path)
        
        # Create validation payload
        validation_payload = {
            "idea_name": extracted_info["idea_name"],
            "idea_concept": extracted_info["idea_concept"]
        }
        
        if custom_weights:
            validation_payload["custom_weights"] = custom_weights
        
        return extracted_info, validation_payload


def process_pitch_deck_file(file_path: str) -> Dict[str, str]:
    """
    Convenience function to process a pitch deck file
    
    Args:
        file_path: Path to the pitch deck file (.pdf, .ppt, or .pptx)
        
    Returns:
        Dictionary with 'idea_name' and 'idea_concept'
        
    Example:
        >>> result = process_pitch_deck_file("my_pitch_deck.pdf")
        >>> print(result['idea_name'])
        >>> print(result['idea_concept'])
    """
    processor = PitchDeckProcessor()
    return processor.process_pitch_deck(file_path)


if __name__ == "__main__":
    # Test the processor
    import sys
    
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
        print(f"Processing pitch deck: {file_path}")
        print("=" * 80)
        
        try:
            processor = PitchDeckProcessor()
            result = processor.process_pitch_deck(file_path)
            
            print(f"\n✅ EXTRACTED INFORMATION:")
            print(f"\n📌 Idea Name: {result['idea_name']}")
            print(f"\n📝 Idea Concept:")
            print(result['idea_concept'])
            print("\n" + "=" * 80)
            
        except Exception as e:
            print(f"❌ Error: {e}")
    else:
        print("Usage: python pitch_deck_processor.py <path_to_pitch_deck.pdf|pptx>")

